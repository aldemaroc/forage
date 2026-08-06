"""Document extraction for Forage: PDF (text layer) and office documents.

Firecrawl-style ``pdf`` / ``document`` engines, Python edition:

- PDF with a text layer  -> pymupdf (fast, no OCR for now)
- .docx                  -> python-docx (paragraphs + tables)
- .xlsx                  -> openpyxl (cell values per sheet)
- .pptx                  -> python-pptx (slide shapes text)
- .rtf                   -> striprtf (pure python)

Detection is by URL extension and/or Content-Type. When the bytes do not
parse as the expected type (e.g. a .pdf URL that redirects to HTML), the
caller falls back to the normal hybrid flow.
"""

from __future__ import annotations

import io
import logging
import re
from typing import Optional, Tuple
from urllib.parse import urlparse

logger = logging.getLogger(__name__)

DOCUMENT_EXTENSIONS = {
    ".pdf": "pdf",
    ".docx": "document",
    ".xlsx": "document",
    ".pptx": "document",
    ".rtf": "document",
}

# Content-Types that mark a response as a supported document.
DOCUMENT_CONTENT_TYPES = {
    "application/pdf": "pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "document",  # .docx
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "document",  # .xlsx
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": "document",  # .pptx
    "application/rtf": "document",
    "text/rtf": "document",
}


def looks_like_document(url: str, content_type: str = "") -> bool:
    """True when the URL path ends with a known extension or the response
    Content-Type declares a supported document."""
    path = (urlparse(url).path or "").lower()
    for ext in DOCUMENT_EXTENSIONS:
        if path.endswith(ext):
            return True
    ctype = (content_type or "").lower().split(";")[0].strip()
    return ctype in DOCUMENT_CONTENT_TYPES


def _doc_kind(url: str, content_type: str = "") -> Optional[str]:
    """Return the document kind ('pdf' | 'document') for url/ctype, else None."""
    path = (urlparse(url).path or "").lower()
    for ext, kind in DOCUMENT_EXTENSIONS.items():
        if path.endswith(ext):
            return kind
    ctype = (content_type or "").lower().split(";")[0].strip()
    return DOCUMENT_CONTENT_TYPES.get(ctype)


def _title_from_url(url: str) -> str:
    """Derive a title from the URL's final path segment."""
    path = (urlparse(url).path or "").rstrip("/")
    name = path.rsplit("/", 1)[-1] or url
    return re.sub(r"[_-]+", " ", name).strip()


def _extract_pdf(data: bytes) -> Tuple[str, str]:
    import fitz  # pymupdf

    doc = fitz.open(stream=data, filetype="pdf")
    try:
        pages = []
        for page in doc:
            pages.append(page.get_text("text"))
        text = "\n\n".join(pages)
        meta_title = (doc.metadata or {}).get("title", "").strip()
    finally:
        doc.close()
    return text, meta_title


def _extract_docx(data: bytes) -> str:
    import docx  # python-docx

    document = docx.Document(io.BytesIO(data))
    parts: list[str] = []
    for para in document.paragraphs:
        if para.text.strip():
            parts.append(para.text)
    for table in document.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            parts.append(" | ".join(cells))
    return "\n\n".join(parts)


def _extract_xlsx(data: bytes) -> str:
    import openpyxl

    wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    parts: list[str] = []
    try:
        for sheet in wb.worksheets:
            parts.append(f"## {sheet.title}")
            for row in sheet.iter_rows(values_only=True):
                cells = ["" if v is None else str(v).strip() for v in row]
                if any(cells):
                    parts.append(" | ".join(cells))
    finally:
        wb.close()
    return "\n\n".join(parts)


def _extract_pptx(data: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    parts: list[str] = []
    for idx, slide in enumerate(prs.slides, 1):
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = "".join(run.text for run in para.runs).strip()
                    if t:
                        slide_text.append(t)
        if slide_text:
            parts.append(f"## Slide {idx}\n" + "\n".join(slide_text))
    return "\n\n".join(parts)


def _extract_rtf(data: bytes) -> str:
    from striprtf.striprtf import rtf_to_text

    try:
        return rtf_to_text(data.decode("utf-8", errors="replace"))
    except Exception:  # noqa: BLE001 (striprtf is permissive, but be safe)
        logger.warning("RTF parse failed, trying latin-1")
        return rtf_to_text(data.decode("latin-1", errors="replace"))


def extract_document_bytes(
    data: bytes,
    url: str,
    content_type: str = "",
    max_chars: int = 100000,
) -> Tuple[str, str, str]:
    """Extract text from document bytes.

    Returns ``(text, title, method_label)``. ``method_label`` is ``pdf`` or
    ``document`` (Firecrawl-style). Raises ``ValueError`` when the bytes do
    not parse as the expected document type (caller falls back).
    """
    kind = _doc_kind(url, content_type)
    if kind is None:
        raise ValueError("not a supported document")

    meta_title = ""
    if kind == "pdf":
        text, meta_title = _extract_pdf(data)
    else:
        path = (urlparse(url).path or "").lower()
        if path.endswith(".docx"):
            text = _extract_docx(data)
        elif path.endswith(".xlsx"):
            text = _extract_xlsx(data)
        elif path.endswith(".pptx"):
            text = _extract_pptx(data)
        elif path.endswith(".rtf"):
            text = _extract_rtf(data)
        else:
            # Content-Type said office document but the URL has no extension;
            # try the parsers in order until one succeeds.
            text = ""
            for parser in (_extract_docx, _extract_xlsx, _extract_pptx, _extract_rtf):
                try:
                    text = parser(data)
                    if text:
                        break
                except Exception:  # noqa: BLE001
                    continue

    text = (text or "").strip()
    if not text:
        raise ValueError("document produced no text (likely scanned / no text layer)")

    text = text[:max_chars]
    title = meta_title or _title_from_url(url)
    return text, title, kind
